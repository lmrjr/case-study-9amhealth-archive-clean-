"""
Generate a branded 9amHealth slide template -> template_9am.pptx (repo root).
Theme derived from brand-tokens.md (Cream bg, Charcoal text, Sunrise accents,
Noto Serif/Sans as STK Bureau stand-ins). Starter slides mirror the case-study
deliverable structure. Edit in LibreOffice Impress; export PDF for the panel.

Run:  .venv/bin/python python-prog/build_template.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- brand tokens ---
CREAM     = RGBColor(0xFF, 0xFC, 0xF3)
CHARCOAL  = RGBColor(0x21, 0x21, 0x21)
BLUE      = RGBColor(0x80, 0xAE, 0xFF)   # Sunrise top
PINK      = RGBColor(0xF7, 0xBD, 0xE6)   # Sunrise mid
ORANGE    = RGBColor(0xFF, 0xBD, 0x70)   # Sunrise bottom
SERIF     = "Noto Serif"                 # STK Bureau Serif stand-in
SANS      = "Noto Sans"                  # STK Bureau Sans stand-in

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def accent_bar(s, color=BLUE, top=Inches(1.55), left=Inches(0.9),
               width=Inches(2.2), height=Inches(0.09)):
    box = s.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box


def text(s, txt, left, top, width, height, size, font=SANS, color=CHARCOAL,
         bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = txt.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.bold = bold
    return tb


def footer(s, label="9amHealth · Data Analyst Case Study"):
    text(s, label, Inches(0.9), Inches(7.02), Inches(11.5), Inches(0.35),
         10, SANS, CHARCOAL)


# 1 — Title
s = slide()
accent_bar(s, BLUE, top=Inches(2.55))
text(s, "Key Drivers of Weight-Loss\nImprovement", Inches(0.9), Inches(2.7),
     Inches(11), Inches(2), 40, SERIF, CHARCOAL, bold=False)
text(s, "9amHealth member cohort · [your name] · [date]", Inches(0.9),
     Inches(4.5), Inches(11), Inches(0.5), 16, SANS, CHARCOAL)
footer(s)

# 2 — Approach
s = slide()
accent_bar(s, BLUE)
text(s, "Analytical Approach", Inches(0.9), Inches(0.7), Inches(11),
     Inches(0.9), 30, SERIF, CHARCOAL)
text(s,
     "• Business question broken into sub-questions + hypotheses\n"
     "• Datasets: demographics, body-weight, engagement, module completion\n"
     "• Member-level join; weight-loss success = Diff > 0 (First − Last)\n"
     "• Explicit assumptions logged (see appendix)",
     Inches(0.9), Inches(1.9), Inches(11.4), Inches(4), 20, SANS, CHARCOAL)
footer(s)

# 3 — Section divider
s = slide()
accent_bar(s, PINK, top=Inches(3.3), width=Inches(3))
text(s, "Key Findings", Inches(0.9), Inches(3.5), Inches(11), Inches(1.2),
     36, SERIF, CHARCOAL)
footer(s)

# 4 — Content + chart placeholder
s = slide()
accent_bar(s, BLUE)
text(s, "Finding: [headline in one line]", Inches(0.9), Inches(0.7),
     Inches(11.4), Inches(0.9), 28, SERIF, CHARCOAL)
ph = s.shapes.add_shape(1, Inches(0.9), Inches(1.9), Inches(7.3), Inches(4.5))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
ph.line.color.rgb = BLUE; ph.line.width = Pt(1)
ph.text_frame.text = "chart here"
ph.text_frame.paragraphs[0].runs[0].font.name = SANS
ph.text_frame.paragraphs[0].runs[0].font.color.rgb = CHARCOAL
text(s, "So what:\n[business implication +\nrecommended action]",
     Inches(8.5), Inches(1.9), Inches(3.8), Inches(4.5), 18, SANS, CHARCOAL)
footer(s)

# 5 — Recommendations
s = slide()
accent_bar(s, ORANGE)
text(s, "Recommendations for Leadership", Inches(0.9), Inches(0.7),
     Inches(11.4), Inches(0.9), 30, SERIF, CHARCOAL)
text(s,
     "1. [action] — [expected impact]\n"
     "2. [action] — [expected impact]\n"
     "3. [action] — [expected impact]",
     Inches(0.9), Inches(1.9), Inches(11.4), Inches(4), 22, SANS, CHARCOAL)
footer(s)

# 6 — Appendix: assumptions
s = slide()
accent_bar(s, PINK)
text(s, "Appendix · Assumptions & Data Notes", Inches(0.9), Inches(0.7),
     Inches(11.4), Inches(0.9), 26, SERIF, CHARCOAL)
text(s,
     "• Demographics has no age column (briefing said it would)\n"
     "• Weight-loss success defined as Diff > 0 (Diff = First − Last)\n"
     "• Join assumed 1:1: bw_detail.User Id ↔ others.Readable Id\n"
     "• Trailing empty column dropped as export artifact",
     Inches(0.9), Inches(1.9), Inches(11.4), Inches(4), 18, SANS, CHARCOAL)
footer(s)

out = "template_9am.pptx"
prs.save(out)
print("wrote", out, "-", len(prs.slides._sldIdLst), "slides")
