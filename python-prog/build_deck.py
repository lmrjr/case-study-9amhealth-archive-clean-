"""
Build the real 9amHealth case-study deck -> deck_9am.pptx (repo root).
Content from deck-outline.md, charts from figures/ (rendered by r-prog/main.R),
brand tokens from brand-tokens.md (Cream bg / Charcoal text / Sunrise accents).

Slides: 1 title · 2 approach (hypotheses + CONSORT) · 3 weight · 4 engagement ·
5 drivers (forest) · 6 hypothesis scorecard · 7 recommendations · 8 appendix
methods · 9 appendix Table 1 · 10 appendix coefficients · 11 appendix power.

Safety: never clobbers an OPEN deck (LibreOffice lock) and backs up the existing
deck before overwrite. Placeholders ([your name]/[date]) left for Luis.

Run:  .venv/bin/python python-prog/build_deck.py   (after r-prog/main.R has written figures/)
"""
import os
import shutil
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# --- brand tokens (brand-guidelines.md §2) ---
CREAM    = RGBColor(0xFF, 0xFC, 0xF3)
CHARCOAL = RGBColor(0x21, 0x21, 0x21)
BLUE     = RGBColor(0x80, 0xAE, 0xFF)   # Sunrise top
PINK     = RGBColor(0xF7, 0xBD, 0xE6)   # Sunrise mid
ORANGE   = RGBColor(0xFF, 0xBD, 0x70)   # Sunrise bottom
SERIF    = "Fraunces"                   # STK Bureau Serif stand-in — free editorial serif
SANS     = "Inter"                      # STK Bureau Sans stand-in — clean grotesque

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG  = os.path.join(ROOT, "figures")

# --- Sunrise light-pattern strip (billboard motif, brand guidelines p.68) ---
# strip runs warm-left -> cool-right, pixel-sampled from the p.68 billboard
SUNRISE   = ["#FFBD70", "#F7BDE6", "#80AEFF"]
AFTERNOON = ["#FFFCF3", "#B0F2CE", "#FFF3C7"]   # p.68: Afternoon strip = educational content
STRIP     = os.path.join(FIG, "strip_sunrise.png")
STRIP_AFT = os.path.join(FIG, "strip_afternoon.png")
TITLE_BG  = os.path.join(FIG, "title_sunrise.png")


def _hex(s):
    return tuple(int(s[i:i + 2], 16) for i in (1, 3, 5))


def _lerp(a, b, t):
    return tuple(int(a[j] + (b[j] - a[j]) * t) for j in range(3))


def gradient_strip(path, stops, w=2400, h=48):
    """Horizontal multi-stop gradient PNG — the bottom light-pattern band."""
    rgb = [_hex(s) for s in stops]
    img = Image.new("RGB", (w, h))
    px = img.load()
    seg = w / (len(rgb) - 1)
    for x in range(w):
        k = min(int(x / seg), len(rgb) - 2)
        c = _lerp(rgb[k], rgb[k + 1], (x - k * seg) / seg)
        for y in range(h):
            px[x, y] = c
    img.save(path)


def _vgrad(w, h, stops):
    """Vertical multi-stop gradient; stops = [(pos 0..1, '#hex'), ...]."""
    img = Image.new("RGB", (w, h))
    for y in range(h):
        t = y / (h - 1)
        for (p0, c0), (p1, c1) in zip(stops, stops[1:]):
            if t <= p1:
                f = 0 if p1 == p0 else (t - p0) / (p1 - p0)
                img.paste(Image.new("RGB", (w, 1), _lerp(_hex(c0), _hex(c1), f)),
                          (0, y))
                break
    return img


def title_background(path, w=1600, h=900):
    """'Refresh' primary light pattern (p.60, colors pixel-sampled from the
    guidelines): warm peach top -> pink -> lilac -> periwinkle field, with an
    orange sun arc cresting ~63% height and fading to salmon at the bottom."""
    img = _vgrad(w, h, [(0.0, "#FDBD8C"), (0.40, "#F9BDC4"),
                        (0.62, "#E4BBEA"), (1.0, "#B6B5F4")])
    crest = int(h * 0.63)
    r = 1535 * w // 1600                              # arc meets edges ~88% height
    cx, cy = w // 2, crest + r                        # sun never fully visible (p.63)
    from PIL import ImageDraw
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).ellipse((cx - r, cy - r, cx + r, cy + r), fill=255)
    sun = _vgrad(w, h, [(0.0, "#FFBD70"), (crest / h, "#FFBD70"),
                        (1.0, "#FABDB4")])
    img.paste(sun, (0, 0), mask)
    img.save(path)


for p, fn in ((STRIP, lambda: gradient_strip(STRIP, SUNRISE)),
              (STRIP_AFT, lambda: gradient_strip(STRIP_AFT, AFTERNOON)),
              (TITLE_BG, lambda: title_background(TITLE_BG))):
    if not os.path.exists(p):
        fn()

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, color, left, top, width, height, shape=1):
    box = s.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box


def eyebrow(s, label, color=BLUE):
    """Eyebrow + 1pt ruler header motif (guidelines p.48 hierarchy, p.91 containers).
    Small color chip keeps the per-slide Sunrise coding."""
    rect(s, color, Inches(0.9), Inches(0.56), Inches(0.14), Inches(0.14))
    text(s, label.upper(), Inches(1.14), Inches(0.47), Inches(10), Inches(0.3),
         11, SANS, CHARCOAL, bold=True)
    rect(s, CHARCOAL, Inches(0.9), Inches(0.84), Inches(11.5), Pt(1))


def text(s, txt, left, top, width, height, size, font=SANS, color=CHARCOAL,
         bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.bold = bold
    return tb


def pic(s, name, left, top, width):
    """Place a figure. Missing figure -> loud error (fail early, don't ship a blank slide)."""
    path = os.path.join(FIG, name)
    if not os.path.exists(path):
        raise FileNotFoundError("missing figure %s — run r-prog/main.R first" % name)
    return s.shapes.add_picture(path, left, top, width=width)


def footer(s, n, strip=STRIP):
    # light-pattern strip along the bottom edge (billboard motif, p.68);
    # Sunrise for brand comms, Afternoon for educational (appendix) slides
    s.shapes.add_picture(strip, Inches(0), Inches(6.92),
                         width=Inches(13.333), height=Inches(0.13))
    text(s, "9amHealth · Data Analyst Case Study", Inches(0.9), Inches(7.12),
         Inches(9), Inches(0.3), 9, SANS, CHARCOAL)
    text(s, str(n), Inches(12.6), Inches(7.12), Inches(0.5), Inches(0.3),
         9, SANS, CHARCOAL, align=PP_ALIGN.RIGHT)


def read_ods(path):
    """Rows of first sheet in an .ods (stdlib zip+xml); skips all-blank rows."""
    import zipfile
    import xml.etree.ElementTree as ET
    T = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
    root = ET.fromstring(zipfile.ZipFile(path).read("content.xml"))
    out = []
    for r in root.iter("{%s}table-row" % T):
        cells = []
        for c in r.findall("{%s}table-cell" % T):
            rep = int(c.get("{%s}number-columns-repeated" % T, 1))
            cells += ["".join(c.itertext())] * min(rep, 8)
        if any(cells):
            out.append(cells)
    return out


def set_cell(cell, txt, bold=False, size=12, align=PP_ALIGN.LEFT, fill=CREAM,
             vpad=0.04):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.09)
    cell.margin_top = Inches(vpad); cell.margin_bottom = Inches(vpad)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.name = SANS; r.font.size = Pt(size)
    r.font.color.rgb = CHARCOAL; r.font.bold = bold


def content(label, title, chart, sowhat, accent=BLUE, n=0, chart_w=7.7):
    """Finding slide, billboard motif: eyebrow + ruler, text left, chart flush-right."""
    s = slide()
    eyebrow(s, label, accent)
    text(s, title, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0),
         26, SERIF, CHARCOAL)
    text(s, sowhat, Inches(0.9), Inches(1.95), Inches(3.9), Inches(4.8),
         15, SANS, CHARCOAL)
    pic(s, chart, Inches(13.333 - 0.5 - chart_w), Inches(2.0), Inches(chart_w))
    footer(s, n)
    return s


# 1 — Title: full Sunrise spectrum + rising sun (intro = high-level brand comms,
# p.54/p.57); charcoal text over spectrum is the only legible option (p.50-51)
s = slide()
s.shapes.add_picture(TITLE_BG, Inches(0), Inches(0),
                     width=Inches(13.333), height=Inches(7.5))
text(s, "What drives clinical weight improvement\namong 9amHealth members?",
     Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.6), 36, SERIF, CHARCOAL,
     align=PP_ALIGN.CENTER)
text(s, "Mostly medication + starting weight + time in program.\n"
        "Engagement & education add a real, smaller, modifiable boost.",
     Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.2), 19, SANS, CHARCOAL,
     bold=True, align=PP_ALIGN.CENTER)
text(s, "Cohort n = 825 · Feb–Sep 2025 · [your name] · [date]",
     Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5), 13, SANS, CHARCOAL,
     align=PP_ALIGN.CENTER)

# 2 — Approach & data: hypotheses left, CONSORT table right (native tables; no consort png in pipeline)
s = slide()
eyebrow(s, "Approach", BLUE)
text(s, "Analytical approach: locked population, clinical outcome, two lenses",
     Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), 26, SERIF, CHARCOAL)
text(s,
     "Sub-questions / hypotheses:\n"
     "H1  Heavier at baseline → larger loss (regression to the mean)\n"
     "H2  Medication (GLP-1) drives more loss than coaching alone\n"
     "H3  Engagement (logging volume) adds benefit on top of medication\n"
     "H4  Module completion adds independent benefit\n"
     "H5  Demographics (sex) relate to weight loss\n"
     "H6  Engagement pays off differently by drug class\n\n"
     "Outcome = % weight change · success = ≥5% loss (clinical bar).\n"
     "Method = baseline-adjusted regression + LASSO (parsimony) + HC3 inference.\n"
     "Full cohort description: Table 1 (appendix).",
     Inches(0.9), Inches(1.85), Inches(6.9), Inches(4.9), 14, SANS, CHARCOAL, space=6)
# CONSORT waterfall as a native table (865 → 827 → 825)
consort = [
    ("Step", "n", "dropped"),
    ("Enrolled universe (demographics)", "865", "—"),
    ("Has BW outcome (first + last)", "827", "38"),
    ("Eligible status (active / finished)", "825", "2"),
]
ctab = s.shapes.add_table(len(consort), 3, Inches(8.1), Inches(2.2),
                          Inches(4.4), Inches(2.2)).table
ctab.first_row = False; ctab.horz_banding = False
ctab.columns[0].width = Inches(2.9); ctab.columns[1].width = Inches(0.75)
ctab.columns[2].width = Inches(0.75)
for i, row in enumerate(consort):
    for j, val in enumerate(row):
        set_cell(ctab.cell(i, j), val, bold=(i == 0),
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
text(s, "CONSORT inclusion waterfall", Inches(8.1), Inches(1.85),
     Inches(4.4), Inches(0.3), 12, SANS, CHARCOAL, bold=True)
footer(s, 2)

# 3 — Weight-change patterns (Objective 1)
content(
    "Findings · Weight",
    "Meaningful weight loss is concentrated in GLP-1 members",
    "pct_distribution_by_drug.png",
    "Weight-change patterns:\n"
    "• Overall median 1.3% loss\n"
    "• 31% hit ≥5% (254 / 825)\n\n"
    "≥5% responders by drug class:\n"
    "• Coaching 5.3%\n"
    "• Generic 30.2%\n"
    "• GLP-1 (All) 61.9%\n\n"
    "Coaching-only median = 0.0%",
    accent=PINK, n=3)

# 4 — Engagement patterns (Objective 2)
content(
    "Findings · Engagement",
    "Engagement isn't uniform — it tracks the medicated population",
    "engage_response_comove.png",
    "Engagement patterns:\n"
    "• 825 / 825 have engagement data\n\n"
    "Median volume_rep:\n"
    "• 7 → 50 → 103\n"
    "  (Coaching → Generic → GLP-1)\n"
    "• mod.mean: 0.31 → 0.90 → 2.70\n\n"
    "Rises in the same order as the\n"
    "responder rate — so how much of\n"
    "it is real? (next slide)",
    accent=ORANGE, n=4)

# 5 — Drivers, confound-adjusted (Objective 3) — centerpiece
content(
    "Findings · Drivers",
    "Engagement: ~half confounded, ~half a real, modifiable lever",
    "driver_forest.png",
    "Drivers (confound-adjusted):\n"
    "Dominant (R² = 0.35): GLP-1\n"
    "+3.0 pts, baseline, exposure.\n\n"
    "Engagement: naive +0.068 →\n"
    "adjusted +0.038 / event,\n"
    "HC3 p < 1e-5.\n"
    "~45% confounding, ~55% real.\n\n"
    "Education adds: mod.mean\n"
    "p = .009, breadth p = .011.\n\n"
    "LASSO 'saw nothing' = predictive\n"
    "parsimony, not inference.",
    accent=BLUE, n=5)

# 6 — Hypothesis scorecard (the attention-director; H1–H6 → verdict)
s = slide()
eyebrow(s, "Findings · Scorecard", PINK)
text(s, "Which questions the analysis answers", Inches(0.9), Inches(0.95),
     Inches(11.5), Inches(1.0), 26, SERIF, CHARCOAL)
pic(s, "hypothesis_scorecard.png", Inches(1.97), Inches(1.85), Inches(9.4))
footer(s, 6)

# 7 — Recommendations: secondary-spectrum tinted containers (p.90-91 motif)
s = slide()
eyebrow(s, "Recommendations", ORANGE)
text(s, "Recommendations", Inches(0.9), Inches(0.95), Inches(11.5),
     Inches(1), 30, SERIF, CHARCOAL)
MIDMORN_T, AFTERN_T, GOLDEN_T = (RGBColor(0xFF, 0xEF, 0xF9),
                                 RGBColor(0xFF, 0xF3, 0xC7),
                                 RGBColor(0xFF, 0xDE, 0xB8))
cols = [
    (MIDMORN_T, "Clinical",
     "GLP-1 is the strongest\nweight-loss lever.\n\n"
     "Coaching-only shows little\nchange on the scale — but body\n"
     "weight can miss body-composition\ngains; measure both.\n\n"
     "Prioritize medication access\n& titration."),
    (AFTERN_T, "Product",
     "Engagement (logging, breadth)\n& module completion are real,\n"
     "modifiable secondary levers.\n\n"
     "Invest in the highest-yield —\nbut it's observational; run a\n"
     "causal test (RCT / matched)\nbefore over-investing."),
    (GOLDEN_T, "Business",
     "Segment on drug class for\nresponder economics.\n\n"
     "Forecast by medication mix,\nwith engagement / education\n"
     "as incremental, controllable\nupside."),
]
x = 0.9
for tint, head, body in cols:
    card = rect(s, tint, Inches(x), Inches(1.85), Inches(3.85), Inches(4.7), shape=5)
    card.adjustments[0] = 0.05
    text(s, head, Inches(x + 0.3), Inches(2.15), Inches(3.3), Inches(0.6),
         20, SERIF, CHARCOAL, bold=True)
    text(s, body, Inches(x + 0.3), Inches(2.95), Inches(3.3), Inches(3.4),
         13.5, SANS, CHARCOAL)
    x += 4.1
footer(s, 7)

# 8 — Appendix: methods & limitations (Afternoon strip: educational content, p.68)
s = slide()
eyebrow(s, "Appendix", PINK)
text(s, "Methods & Limitations", Inches(0.9), Inches(0.95),
     Inches(11.5), Inches(1), 26, SERIF, CHARCOAL)
text(s,
     "• Cohort exclusions: 865 → 827 (has weight) → 825 (active / finished)\n"
     "• No age column (briefing said it would exist) — key limitation\n"
     "• ≥5% = clinically meaningful weight-loss threshold\n"
     "• Outcome = % change (decouples baseline), HC3 robust SE (heteroskedastic)\n"
     "• LASSO (prediction / parsimony) vs HC3 (inference) — why they disagree\n"
     "• Multiplicity: Bonferroni α = 0.005 → only volume_rep unambiguous\n"
     "• Observational — associations, not causal; engagement may proxy adherence\n"
     "• Endpoint: body weight can't separate fat vs lean mass (coaching ~0 may be recomposition)\n"
     "• Full cohort description in Table 1 (next slide)",
     Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8), 15, SANS, CHARCOAL,
     space=8)
footer(s, 8, STRIP_AFT)

# 9 — Appendix: Table 1 cohort description (native table from figures/demo-table.ods)
s = slide()
eyebrow(s, "Appendix", BLUE)
text(s, "Table 1 · Cohort description by treatment group", Inches(0.9),
     Inches(0.95), Inches(11.5), Inches(0.7), 22, SERIF, CHARCOAL)
t1 = read_ods(os.path.join(FIG, "demo-table.ods"))
tab = s.shapes.add_table(len(t1), 6, Inches(0.6), Inches(1.42),
                         Inches(12.13), Inches(5.2)).table
tab.first_row = False; tab.horz_banding = False
tab.columns[0].width = Inches(2.83)
for j in range(1, 6):
    tab.columns[j].width = Inches(1.86)
for i, row in enumerate(t1):
    section = i > 0 and row[0] and not any(row[1:6])
    for j in range(6):
        set_cell(tab.cell(i, j), row[j] if j < len(row) else "",
                 bold=(i == 0 or section), size=8 if (i == 0 or section) else 7.5,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT,
                 fill=GOLDEN_T if (i == 0 or section) else CREAM, vpad=0.015)
    tab.rows[i].height = Inches(0.15)
footer(s, 9, STRIP_AFT)

# 10 — Appendix: locked-model coefficient table
s = slide()
eyebrow(s, "Appendix", BLUE)
text(s, "Locked-model coefficients (HC3)", Inches(0.9), Inches(0.95),
     Inches(11.5), Inches(1), 26, SERIF, CHARCOAL)
pic(s, "coef_table.png", Inches(1.4), Inches(1.95), Inches(10.5))
footer(s, 10, STRIP_AFT)

# 11 — Appendix: power & precision
s = slide()
eyebrow(s, "Appendix", ORANGE)
text(s, "Power & precision", Inches(0.9), Inches(0.95),
     Inches(11.5), Inches(1), 26, SERIF, CHARCOAL)
pic(s, "power_precision.png", Inches(1.0), Inches(1.95), Inches(11.3))
footer(s, 11, STRIP_AFT)

# --- save: never clobber an OPEN deck; back up an existing one first ---
out  = os.path.join(ROOT, "deck_9am.pptx")
lock = os.path.join(ROOT, ".~lock.deck_9am.pptx#")
if os.path.exists(lock):
    out = os.path.join(ROOT, "deck_9am_rebuilt.pptx")
    print("deck_9am.pptx is OPEN (lock present) — writing", os.path.basename(out),
          "instead. Close LibreOffice and rerun to overwrite the live deck.")
elif os.path.exists(out):
    bak = out.replace(".pptx", time.strftime("-%Y%m%d-%H%M%S.pptx.bak"))
    shutil.copy2(out, bak)
    print("backed up existing deck ->", os.path.basename(bak))
prs.save(out)
print("wrote", out, "-", len(prs.slides._sldIdLst), "slides")
